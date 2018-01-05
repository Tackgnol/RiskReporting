from pptx import Presentation

class PowerPoint():
    def __init__(self, pathToPPT = None):
        try:
            self.Presentation = Presentation(pathToPPT)
        except:
            print('Failed to open the presentation creating a new one')
            self.Presentation = Presentation()

    def AddSlide(self, layoutName):
        slideLayout = self.Presentation.slide_layouts[layoutName]
        slide = self.Presentation.slides.add_slide(slideLayout)




